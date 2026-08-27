# -*- coding: utf-8 -*-
"""manage дотоод систем — сургалт (яг 4 слайд, 16:9)."""

from __future__ import annotations

import struct
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

NAVY = RGBColor(0x0F, 0x2F, 0x63)
NAVY2 = RGBColor(0x1C, 0x55, 0xA5)
NAVY700 = RGBColor(0x1E, 0x40, 0x6B)
NAVY_SOFT = RGBColor(0xE8, 0xEE, 0xF6)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x47, 0x55, 0x69)
DARK = RGBColor(0x0F, 0x17, 0x2A)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x64, 0x74, 0x8B)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
SKY = RGBColor(0x0E, 0xA5, 0xE9)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)

W = Inches(13.333)
H = Inches(7.5)
TOTAL = 4

ROOT = Path(__file__).resolve().parent
LOGIN_PNG = ROOT / "slides" / "login.png"
EMBLEM_PNG = ROOT.parent / "public" / "images" / "emblem.png"
OUT = ROOT / "manage-dotood-sistem-surgalt.pptx"


def set_run(run, size=14, bold=False, color=DARK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font)


def add_text_box(slide, l, t, w, h, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_lines(slide, l, t, w, h, lines, size=13, color=DARK, spacing=6, bold=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        if isinstance(item, tuple):
            text, is_bold, col = item
        else:
            text, is_bold, col = item, bold, color
        run = p.add_run()
        run.text = text
        set_run(run, size=size, bold=is_bold, color=col)
    return box


def add_rect(slide, l, t, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    return shape


def add_round(slide, l, t, w, h, fill, line=None, adj=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.adjustments[0] = adj
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.9)
    return shape


def footer(slide, page):
    add_rect(slide, 0, Inches(7.22), W, Inches(0.28), NAVY)
    add_text_box(
        slide, Inches(0.35), Inches(7.22), Inches(10.5), Inches(0.28),
        "manage дотоод систем  ·  Дорноговь аймгийн ЗДТГ  ·  manage.dornogovi.gov.mn",
        size=11, color=WHITE,
    )
    add_text_box(
        slide, Inches(11.5), Inches(7.22), Inches(1.5), Inches(0.28),
        f"{page} / {TOTAL}",
        size=11, color=WHITE, align=PP_ALIGN.RIGHT,
    )


def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(1.08), NAVY)
    add_rect(slide, 0, Inches(1.08), Inches(0.11), Inches(6.14), ORANGE)
    add_text_box(slide, Inches(0.45), Inches(0.18), Inches(12.4), Inches(0.48), title, size=24, bold=True, color=WHITE)
    if subtitle:
        add_text_box(
            slide, Inches(0.45), Inches(0.64), Inches(12.4), Inches(0.34),
            subtitle, size=13, color=RGBColor(0xBF, 0xDB, 0xFE),
        )


def png_size(path: Path):
    with path.open("rb") as f:
        sig = f.read(8)
        if sig != b"\x89PNG\r\n\x1a\n":
            return None
        length, chunk = struct.unpack(">I4s", f.read(8))
        if chunk != b"IHDR" or length < 8:
            return None
        w, h = struct.unpack(">II", f.read(8))
        return w, h


def fit_picture(slide, path, l, t, max_w, max_h):
    size = png_size(path)
    if not size:
        return slide.shapes.add_picture(str(path), l, t, width=max_w)
    pw, ph = size
    scale = min(max_w / pw, max_h / ph)
    w = int(pw * scale)
    h = int(ph * scale)
    x = int(l + (max_w - w) / 2)
    y = int(t + (max_h - h) / 2)
    return slide.shapes.add_picture(str(path), Emu(x), Emu(y), width=Emu(w), height=Emu(h))


def maybe_emblem(slide, l, t, size):
    if EMBLEM_PNG.exists():
        slide.shapes.add_picture(str(EMBLEM_PNG), l, t, width=size, height=size)


def slide_login(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, W, H, LIGHT)
    add_rect(s, 0, 0, W, Inches(1.48), NAVY)
    add_rect(s, 0, Inches(1.48), W, Inches(0.08), ORANGE)

    maybe_emblem(s, Inches(0.4), Inches(0.26), Inches(0.92))
    add_text_box(
        s, Inches(1.5), Inches(0.18), Inches(11.4), Inches(0.5),
        "Дорноговь аймгийн ЗДТГ — Дотоод нэгдсэн систем",
        size=24, bold=True, color=WHITE,
    )
    add_text_box(
        s, Inches(1.5), Inches(0.7), Inches(11.4), Inches(0.32),
        "https://manage.dornogovi.gov.mn",
        size=16, bold=True, color=RGBColor(0xFD, 0xBA, 0x74),
    )
    add_text_box(
        s, Inches(1.5), Inches(1.04), Inches(11.4), Inches(0.3),
        "Сургалт  ·  нэвтрэх боломжууд  ·  бодит нэвтрэх хуудас",
        size=13, color=RGBColor(0xBF, 0xDB, 0xFE),
    )

    add_round(s, Inches(0.32), Inches(1.74), Inches(4.55), Inches(5.28), WHITE, LINE, adj=0.04)
    add_text_box(
        s, Inches(0.48), Inches(1.82), Inches(4.25), Inches(0.28),
        "Харагдах байдал — нэвтрэх", size=12, bold=True, color=NAVY,
    )
    if LOGIN_PNG.exists():
        fit_picture(s, LOGIN_PNG, Inches(0.55), Inches(2.12), Inches(4.1), Inches(4.72))
    else:
        add_text_box(s, Inches(0.5), Inches(4.0), Inches(4.2), Inches(0.4), "login.png олдсонгүй", size=13, color=MUTED, align=PP_ALIGN.CENTER)

    add_text_box(s, Inches(5.1), Inches(1.74), Inches(7.85), Inches(0.34), "Нэвтрэх боломжууд — 3 арга", size=16, bold=True, color=NAVY)

    methods = [
        ("1", "Утас + нууц үг", "Үндсэн арга. 8 оронтой дугаар, нууц үг, «Нэвтрэх»."),
        ("2", "QR кодоор нэвтрэх", "Компьютер дээр код гарна. Утаснаасаа уншуулаад зөвшөөрнө."),
        ("3", "И-мэйлээр нэвтрэх", "Нөөц арга. «И-мэйлээр нэвтрэх» дарж хаяг + нууц үг оруулна."),
    ]
    y = Inches(2.14)
    for num, title, body in methods:
        add_round(s, Inches(5.1), y, Inches(7.85), Inches(0.92), WHITE, LINE, adj=0.08)
        add_round(s, Inches(5.28), y + Inches(0.22), Inches(0.48), Inches(0.48), NAVY, adj=0.2)
        add_text_box(s, Inches(5.28), y + Inches(0.28), Inches(0.48), Inches(0.38), num, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, Inches(5.95), y + Inches(0.14), Inches(6.8), Inches(0.32), title, size=15, bold=True, color=NAVY)
        add_text_box(s, Inches(5.95), y + Inches(0.48), Inches(6.8), Inches(0.36), body, size=13, color=SLATE)
        y += Inches(1.02)

    add_round(s, Inches(5.1), Inches(5.24), Inches(7.85), Inches(0.82), RGBColor(0xFF, 0xF7, 0xED), RGBColor(0xFD, 0xBA, 0x74), adj=0.06)
    add_text_box(s, Inches(5.28), Inches(5.32), Inches(7.5), Inches(0.28), "Нууц үг", size=13, bold=True, color=ORANGE)
    add_text_box(
        s, Inches(5.28), Inches(5.6), Inches(7.5), Inches(0.38),
        "Утасны сүүлийн 4 орон + нэрийн латин бичлэг.   Жишээ: А.Номин / 99178904  →  8904Nomin",
        size=13, color=DARK,
    )

    add_round(s, Inches(5.1), Inches(6.16), Inches(7.85), Inches(0.82), WHITE, LINE, adj=0.06)
    add_text_box(
        s, Inches(5.28), Inches(6.28), Inches(7.5), Inches(0.58),
        "«Намайг сана» — энэ төхөөрөмж дээр дараагийн удаа нэрийг санана.\n"
        "И-мэйл: nomin@dornogovi.gov.mn хэлбэртэй. Нэвтрэх нэр нь үндсэндээ утас.",
        size=13, color=SLATE,
    )

    footer(s, 1)


def _sidebar_item(slide, x, y, w, label, active=False):
    bg = NAVY if active else WHITE
    fg = WHITE if active else DARK
    add_round(slide, x, y, w, Inches(0.28), bg, adj=0.12)
    add_text_box(slide, x + Inches(0.12), y + Inches(0.02), w - Inches(0.16), Inches(0.24), label, size=10, bold=active, color=fg)


def slide_chrome(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, W, H, LIGHT)
    header_bar(s, "Харагдах байдал", "Зүүн цэс · толгой · ажлын хэсэг   ·   компьютер / утас")

    frame_l, frame_t = Inches(0.32), Inches(1.26)
    frame_w, frame_h = Inches(9.05), Inches(5.76)
    add_round(s, frame_l, frame_t, frame_w, frame_h, RGBColor(0xCB, 0xD5, 0xE1), adj=0.025)
    inner_l, inner_t = frame_l + Inches(0.05), frame_t + Inches(0.05)
    inner_w, inner_h = frame_w - Inches(0.1), frame_h - Inches(0.1)

    side_w = Inches(2.62)
    add_rect(s, inner_l, inner_t, side_w, inner_h, WHITE)
    add_rect(s, inner_l + side_w, inner_t, Inches(0.012), inner_h, LINE)

    maybe_emblem(s, inner_l + Inches(0.1), inner_t + Inches(0.1), Inches(0.36))
    add_text_box(s, inner_l + Inches(0.52), inner_t + Inches(0.08), Inches(2.0), Inches(0.22), "manage", size=12, bold=True, color=NAVY)
    add_text_box(s, inner_l + Inches(0.52), inner_t + Inches(0.28), Inches(2.0), Inches(0.18), "дотоод систем", size=9, color=MUTED)
    add_rect(s, inner_l, inner_t + Inches(0.5), side_w, Inches(0.01), LINE)

    groups = [
        ("САМБАР", [("Албан хаагчийн самбар", True)]),
        ("ХОЛБООСОН СИСТЕМҮҮД", [("Төрийн ERP", False), ("Шилэн данс", False)]),
        ("АЖЛЫН УДИРДЛАГА", [("Үүрэг даалгавар", False), ("Ажлын хэсэг", False), ("Төлөвлөгөө", False)]),
        ("БИЧИГ ХЭРЭГ", [("Захирамж, тушаал", False)]),
        ("ХҮНИЙ НӨӨЦ", [("Утасны жагсаалт", False), ("Чөлөөний бүртгэл", False)]),
        ("МЭДЭЭЛЭЛ, СУРГАЛТ", [("Гарын авлага, сургалт", False)]),
    ]
    y = inner_t + Inches(0.56)
    ix = inner_l + Inches(0.08)
    iw = side_w - Inches(0.16)
    for gname, items in groups:
        add_text_box(s, ix, y, iw, Inches(0.16), gname, size=7, bold=True, color=MUTED)
        y += Inches(0.16)
        for label, active in items:
            _sidebar_item(s, ix, y, iw, label, active=active)
            y += Inches(0.3)
        y += Inches(0.02)

    add_rect(s, inner_l, inner_t + inner_h - Inches(0.62), side_w, Inches(0.01), LINE)
    add_text_box(s, ix, inner_t + inner_h - Inches(0.56), iw, Inches(0.2), "Сан", size=10, color=SLATE)
    add_text_box(s, ix, inner_t + inner_h - Inches(0.36), iw, Inches(0.2), "Хандах эрх", size=10, color=SLATE)
    add_text_box(s, ix, inner_t + inner_h - Inches(0.18), iw, Inches(0.16), "Системийн тохиргоо", size=9, color=MUTED)

    main_l = inner_l + side_w
    main_w = inner_w - side_w
    add_rect(s, main_l, inner_t, main_w, Inches(0.5), WHITE)
    add_rect(s, main_l, inner_t + Inches(0.5), main_w, Inches(0.01), LINE)
    add_text_box(s, main_l + Inches(0.16), inner_t + Inches(0.1), Inches(4.0), Inches(0.3), "Албан хаагчийн самбар", size=13, bold=True, color=NAVY)

    ux = main_l + main_w - Inches(1.78)
    add_round(s, ux, inner_t + Inches(0.09), Inches(1.62), Inches(0.32), LIGHT, LINE, adj=0.2)
    add_round(s, ux + Inches(0.06), inner_t + Inches(0.13), Inches(0.24), Inches(0.24), NAVY, adj=0.5)
    add_text_box(s, ux + Inches(0.34), inner_t + Inches(0.11), Inches(1.22), Inches(0.26), "А.Номин", size=10, bold=True, color=NAVY)

    add_rect(s, main_l, inner_t + Inches(0.51), main_w, inner_h - Inches(0.51), RGBColor(0xF8, 0xFA, 0xFC))
    add_text_box(s, main_l + Inches(0.18), inner_t + Inches(0.6), Inches(5.9), Inches(0.26), "Засаг даргын Тамгын газар", size=13, bold=True, color=NAVY)
    add_text_box(s, main_l + Inches(0.18), inner_t + Inches(0.86), Inches(5.9), Inches(0.22), "Албан хаагчийн товч үзүүлэлт, явц.", size=11, color=MUTED)

    card_specs = [
        ("Чөлөө", "3", RGBColor(0xFF, 0xFB, 0xEB), RGBColor(0xB4, 0x53, 0x09), AMBER),
        ("Томилолт", "1", RGBColor(0xF0, 0xF9, 0xFF), RGBColor(0x03, 0x69, 0xA1), SKY),
        ("Төлөвлөгөө", "2", RGBColor(0xEC, 0xFD, 0xF5), RGBColor(0x04, 0x78, 0x57), EMERALD),
        ("Ажлын хэсэг", "4", RGBColor(0xF5, 0xF3, 0xFF), RGBColor(0x6D, 0x28, 0xD9), VIOLET),
        ("Үүрэг %", "68%", RGBColor(0xF8, 0xFA, 0xFC), RGBColor(0x33, 0x41, 0x55), MUTED),
    ]
    cw = Inches(1.14)
    for i, (lab, val, bg, fg, dot) in enumerate(card_specs):
        cx = main_l + Inches(0.16) + i * Inches(1.22)
        cy = inner_t + Inches(1.18)
        add_round(s, cx, cy, cw, Inches(0.92), bg, LINE, adj=0.1)
        add_round(s, cx + cw - Inches(0.18), cy + Inches(0.1), Inches(0.1), Inches(0.1), dot, adj=0.5)
        add_text_box(s, cx + Inches(0.08), cy + Inches(0.1), cw - Inches(0.16), Inches(0.24), lab, size=9, bold=True, color=fg)
        add_text_box(s, cx + Inches(0.08), cy + Inches(0.4), cw - Inches(0.16), Inches(0.4), val, size=16, bold=True, color=fg)

    add_round(s, main_l + Inches(0.16), inner_t + Inches(2.26), main_w - Inches(0.32), Inches(3.1), WHITE, LINE, adj=0.04)
    add_text_box(s, main_l + Inches(0.32), inner_t + Inches(2.4), main_w - Inches(0.6), Inches(0.28), "Зүүн цэсээр бүх хуудас руу орно", size=13, bold=True, color=NAVY)
    add_lines(s, main_l + Inches(0.32), inner_t + Inches(2.74), main_w - Inches(0.6), Inches(2.4), [
        "Цэс эрхээс хамаарна — хаасан цэс харагдахгүй.",
        "Толгой: хуудасны нэр · QR уншигч · мэдэгдэл · апп суулгах · хэрэглэгч.",
        "Цэсийг компьютераар хурааж болно.",
        "«Үүрэг даалгавар» — Ажлын удирдлага бүлэгт.",
        "Эхний хуудас — Албан хаагчийн самбар.",
    ], size=12, color=SLATE, spacing=6)

    add_round(s, Inches(9.55), Inches(1.26), Inches(3.42), Inches(2.7), WHITE, LINE, adj=0.05)
    add_rect(s, Inches(9.55), Inches(1.26), Inches(0.1), Inches(2.7), NAVY)
    add_text_box(s, Inches(9.8), Inches(1.38), Inches(3.0), Inches(0.3), "Компьютер", size=14, bold=True, color=NAVY)
    add_lines(s, Inches(9.8), Inches(1.74), Inches(3.0), Inches(2.05), [
        "Зүүн цэс байнга харагдана.",
        "Дараалал: Самбар → Холбосон системүүд → Ажлын удирдлага → Бичиг хэрэг → Хүний нөөц.",
        "Доор: Сан, Хандах эрх, тохиргоо.",
    ], size=12, color=SLATE, spacing=5)

    add_round(s, Inches(9.55), Inches(4.12), Inches(3.42), Inches(2.9), WHITE, LINE, adj=0.05)
    add_rect(s, Inches(9.55), Inches(4.12), Inches(0.1), Inches(2.9), ORANGE)
    add_text_box(s, Inches(9.8), Inches(4.24), Inches(3.0), Inches(0.3), "Утас / PWA", size=14, bold=True, color=NAVY)
    add_lines(s, Inches(9.8), Inches(4.58), Inches(3.0), Inches(2.25), [
        "Цэс товчоор зүүн цэс нээгдэнэ.",
        "Android: Chrome → апп болгон суулгана.",
        "iPhone: Safari → «Нүүр дэлгэцэд нэмэх».",
        "Апп горимд хөтчийн хаяг харагдахгүй.",
    ], size=12, color=SLATE, spacing=5)

    footer(s, 2)


def _dash_card(slide, l, t, w, h, label, value, hint, bg, fg, dot):
    add_round(slide, l, t, w, h, bg, RGBColor(0xD6, 0xD3, 0xD1), adj=0.08)
    add_round(slide, l + w - Inches(0.28), t + Inches(0.14), Inches(0.14), Inches(0.14), dot, adj=0.5)
    add_text_box(slide, l + Inches(0.12), t + Inches(0.1), w - Inches(0.4), Inches(0.4), label, size=11, bold=True, color=fg)
    add_text_box(slide, l + Inches(0.12), t + Inches(0.5), w - Inches(0.24), Inches(0.4), value, size=22, bold=True, color=fg)
    if hint:
        add_text_box(slide, l + Inches(0.12), t + Inches(0.96), w - Inches(0.24), Inches(0.26), hint, size=10, color=fg)


def slide_dashboard(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, W, H, LIGHT)
    header_bar(s, "Албан хаагчийн самбар", "Нэвтэрсний дараах эхний хуудас  ·  зүүн цэс: Самбар → Албан хаагчийн самбар")

    cards = [
        ("Хүлээгдэж буй чөлөө", "3", "Дэлгэрэнгүй харах", RGBColor(0xFF, 0xFB, 0xEB), RGBColor(0xB4, 0x53, 0x09), AMBER),
        ("Идэвхтэй томилолт", "1", "Дэлгэрэнгүй харах", RGBColor(0xF0, 0xF9, 0xFF), RGBColor(0x03, 0x69, 0xA1), SKY),
        ("Идэвхтэй төлөвлөгөө", "2", "Дэлгэрэнгүй харах", RGBColor(0xEC, 0xFD, 0xF5), RGBColor(0x04, 0x78, 0x57), EMERALD),
        ("Ажлын хэсэг", "4", "Дэлгэрэнгүй харах", RGBColor(0xF5, 0xF3, 0xFF), RGBColor(0x6D, 0x28, 0xD9), VIOLET),
        ("Үүргийн дундаж", "68%", "", RGBColor(0xF8, 0xFA, 0xFC), RGBColor(0x33, 0x41, 0x55), MUTED),
    ]
    gap = Inches(0.14)
    card_w = Inches(2.42)
    x0 = Inches(0.4)
    for i, (lab, val, hint, bg, fg, dot) in enumerate(cards):
        _dash_card(s, x0 + i * (card_w + gap), Inches(1.28), card_w, Inches(1.32), lab, val, hint, bg, fg, dot)

    add_round(s, Inches(0.4), Inches(2.78), Inches(8.4), Inches(4.22), WHITE, LINE, adj=0.04)
    add_text_box(s, Inches(0.58), Inches(2.9), Inches(6.2), Inches(0.32), "Сүүлийн чөлөө", size=15, bold=True, color=NAVY)
    add_text_box(s, Inches(7.15), Inches(2.92), Inches(1.4), Inches(0.3), "Бүгд", size=13, bold=True, color=NAVY2, align=PP_ALIGN.RIGHT)

    rows = [
        ("Б.Бат", "Ээлжийн амралт", "хүлээгдэж буй"),
        ("Д.Саран", "Өвчний чөлөө", "хүлээгдэж буй"),
        ("Г.Эрдэнэ", "Цалинтай чөлөө", "батлагдсан"),
    ]
    y = Inches(3.36)
    for name, typ, st in rows:
        add_rect(s, Inches(0.58), y + Inches(0.5), Inches(8.04), Inches(0.01), LINE)
        add_text_box(s, Inches(0.58), y, Inches(5.5), Inches(0.3), f"{name}  ·  {typ}", size=13, color=DARK)
        add_text_box(s, Inches(6.2), y, Inches(2.4), Inches(0.3), st, size=12, color=MUTED, align=PP_ALIGN.RIGHT)
        y += Inches(0.56)

    add_round(s, Inches(0.58), Inches(5.2), Inches(8.04), Inches(1.58), RGBColor(0xF8, 0xFA, 0xFC), LINE, adj=0.06)
    add_lines(s, Inches(0.74), Inches(5.32), Inches(7.72), Inches(1.36), [
        "1.  Карт дээр дарна → доор сүүлийн бүртгэл нээгдэнэ.",
        "2.  «Бүгд» → тухайн бүртгэлийн хуудас руу орно.",
        "3.  Дахин дарахад самбар хаагдана.",
        "4.  Үүргийн дундаж нээгдэхгүй — зөвхөн хувь харуулна.",
    ], size=13, color=SLATE, spacing=4)

    add_round(s, Inches(8.98), Inches(2.78), Inches(3.95), Inches(4.22), WHITE, LINE, adj=0.04)
    add_rect(s, Inches(8.98), Inches(2.78), Inches(0.1), Inches(4.22), ORANGE)
    add_text_box(s, Inches(9.24), Inches(2.94), Inches(3.5), Inches(0.32), "Таван карт", size=15, bold=True, color=NAVY)
    add_lines(s, Inches(9.24), Inches(3.36), Inches(3.5), Inches(3.4), [
        "Чөлөө — хүлээгдэж буй тоо.",
        "Томилолт — идэвхтэй явц.",
        "Төлөвлөгөө — нээлттэй төлөвлөгөө.",
        "Ажлын хэсэг — багийн тоо.",
        "Үүргийн дундаж — бүх үүргийн %.",
        "",
        "Өдөр тутмын тоймоо эндээс авна.",
    ], size=13, color=SLATE, spacing=6)

    footer(s, 3)


def slide_tasks(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, W, H, LIGHT)
    header_bar(s, "Үүрэг даалгавар", "Зүүн цэс: Ажлын удирдлага → Үүрэг даалгавар   ·   нүд дээр дарж шууд засна")

    actions = [
        ("Word оруулах", NAVY, WHITE, Inches(1.85)),
        ("Татах  ▾", WHITE, NAVY, Inches(1.45)),
        ("Үүрэг чиглэл нэмэх", ORANGE, WHITE, Inches(2.2)),
    ]
    ax = Inches(0.4)
    for label, bg, fg, bw in actions:
        add_round(s, ax, Inches(1.24), bw, Inches(0.36), bg, LINE if bg == WHITE else None, adj=0.15)
        add_text_box(s, ax, Inches(1.27), bw, Inches(0.3), label, size=11, bold=True, color=fg, align=PP_ALIGN.CENTER)
        ax += bw + Inches(0.1)
    add_text_box(s, Inches(8.4), Inches(1.26), Inches(4.5), Inches(0.32), "Татах: Word · Excel · PDF", size=12, color=MUTED)

    add_round(s, Inches(0.4), Inches(1.7), Inches(12.55), Inches(0.46), WHITE, LINE, adj=0.08)
    add_round(s, Inches(0.48), Inches(1.76), Inches(2.2), Inches(0.34), NAVY, adj=0.12)
    add_text_box(s, Inches(0.48), Inches(1.8), Inches(2.2), Inches(0.28), "Үүрэг чиглэл", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(2.82), Inches(1.8), Inches(4.7), Inches(0.28), "Бэлтгэл ажил хангах төлөвлөгөө", size=12, color=SLATE)
    add_text_box(s, Inches(7.7), Inches(1.8), Inches(2.4), Inches(0.28), "+ Хэсэг нэмэх", size=12, bold=True, color=NAVY2)

    add_round(s, Inches(0.4), Inches(2.28), Inches(12.55), Inches(0.7), WHITE, LINE, adj=0.06)
    add_round(s, Inches(0.55), Inches(2.4), Inches(0.44), Inches(0.44), NAVY_SOFT, adj=0.15)
    add_text_box(s, Inches(1.1), Inches(2.34), Inches(4.4), Inches(0.26), "Хэрэгжилтийн дашбоард", size=13, bold=True, color=NAVY)
    add_text_box(s, Inches(1.1), Inches(2.58), Inches(4.4), Inches(0.28), "Дарж дэлгэрэнгүй графикаар харна    Нийт 12   Дууссан 5   Эхлээгүй 3", size=11, color=SLATE)
    add_text_box(s, Inches(11.2), Inches(2.38), Inches(1.55), Inches(0.42), "62%", size=22, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
    add_round(s, Inches(5.7), Inches(2.56), Inches(5.35), Inches(0.14), RGBColor(0xE2, 0xE8, 0xF0), adj=0.5)
    add_round(s, Inches(5.7), Inches(2.56), Inches(3.3), Inches(0.14), EMERALD, adj=0.5)

    cols = [
        (Inches(0.5), "№"),
        (Inches(3.4), "Үүрэг чиглэл"),
        (Inches(2.3), "Хариуцах эзэн"),
        (Inches(2.7), "Хяналт тавих албан тушаалтан"),
        (Inches(2.15), "Хэрэгжилт"),
        (Inches(1.5), "Биелэлтийн хувь"),
    ]
    tx, ty = Inches(0.4), Inches(3.12)
    th = Inches(0.34)
    x = tx
    for w, lab in cols:
        add_rect(s, x, ty, w, th, NAVY_SOFT, LINE)
        add_text_box(s, x + Inches(0.05), ty + Inches(0.04), w - Inches(0.08), Inches(0.26), lab, size=9, bold=True, color=NAVY700)
        x += w

    data = [
        ("1", "Зудын эсрэг бэлтгэл хангах", "А.Номин", "Хэлтсийн дарга", "Хурлын шийдвэр гарсан", "80%"),
        ("2", "Сумдын тайлан нэгтгэх", "А.Бадрал / Ц.Саран", "Дэд дарга", "Цуглуулж байна", "45%"),
        ("3", "Сургалтын хуваарь батлуулах", "Б.Эрдэнэ", "Хэлтсийн дарга", "Эхлээгүй", "0%"),
    ]
    row_h = Inches(0.4)
    for r, row in enumerate(data):
        y = ty + th + r * row_h
        bg = WHITE if r % 2 == 0 else LIGHT
        x = tx
        for i, (w, _) in enumerate(cols):
            add_rect(s, x, y, w, row_h, bg, LINE)
            fg = EMERALD if i == 5 and row[i] not in ("0%", "45%") else (ORANGE if row[i] == "0%" else DARK)
            add_text_box(s, x + Inches(0.05), y + Inches(0.06), w - Inches(0.08), Inches(0.28), row[i], size=11, bold=(i in (0, 5)), color=fg)
            x += w

    tips = [
        ("Нүд засах", "Нүд дээр дарж шууд бичнэ. Тусдаа «Хадгалах» байхгүй."),
        ("Хүн сонгох", "Хариуцагч / хяналтыг утасны жагсаалтаас нэрээр сонгоно."),
        ("Хэсэг нэмэх", "«+ Хэсэг нэмэх» — шинэ таб. Загвар: Үүрэг чиглэл эсвэл Бэлтгэл ажил."),
    ]
    tw = Inches(4.1)
    for i, (title, body) in enumerate(tips):
        x = Inches(0.4) + i * (tw + Inches(0.12))
        add_round(s, x, Inches(5.18), tw, Inches(1.82), WHITE, LINE, adj=0.06)
        add_rect(s, x, Inches(5.18), Inches(0.1), Inches(1.82), ORANGE if i == 2 else NAVY)
        add_text_box(s, x + Inches(0.22), Inches(5.28), tw - Inches(0.32), Inches(0.32), title, size=14, bold=True, color=NAVY)
        add_text_box(s, x + Inches(0.22), Inches(5.62), tw - Inches(0.32), Inches(1.2), body, size=13, color=SLATE)

    footer(s, 4)


def build(out: Path | None = None):
    path = out or OUT
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slide_login(prs)
    slide_chrome(prs)
    slide_dashboard(prs)
    slide_tasks(prs)
    prs.save(path)
    return path


def verify(path: Path):
    if not path.exists() or path.stat().st_size < 10_000:
        raise SystemExit(f"PPTX missing or too small: {path}")
    with zipfile.ZipFile(path) as z:
        names = z.namelist()
        if "[Content_Types].xml" not in names or "ppt/presentation.xml" not in names:
            raise SystemExit("Not a valid PPTX zip")
        media = [n for n in names if n.startswith("ppt/media/")]
        if LOGIN_PNG.exists() and not media:
            raise SystemExit("Expected embedded screenshot in PPTX")
    prs = Presentation(str(path))
    n = len(prs.slides)
    if n != TOTAL:
        raise SystemExit(f"Expected {TOTAL} slides, got {n}")
    if prs.slide_width != W or prs.slide_height != H:
        raise SystemExit("Slide size is not 16:9 widescreen")
    print(f"OK  {path}  slides={n}  size={path.stat().st_size}  login_png={LOGIN_PNG.exists()}")


if __name__ == "__main__":
    target = OUT
    try:
        out = build(target)
    except PermissionError:
        target = ROOT / "manage-dotood-sistem-surgalt.new.pptx"
        out = build(target)
        print(f"LOCKED — saved as {out}")
    verify(out)
    print(out)
